"""PPTX text extraction from slide shapes, tables, and speaker notes."""

from pptx import Presentation


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from all slides: titles, body text, tables, and notes."""
    presentation = Presentation(file_path)
    slides_text: list[str] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = [f"[Slide {slide_number}]"]

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip(" |"):
                        parts.append(row_text)

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"Notes: {notes}")

        slides_text.append("\n".join(parts))

    return "\n\n".join(slides_text)
